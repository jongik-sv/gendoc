#!/usr/bin/env python3
"""
PPTX 병합 및 LLM 템플릿 변환 스크립트

사용법:
    # 분석
    python pptx_merge.py --analyze PPT기본양식.pptx

    # 병합 (템플릿 변환 포함)
    python pptx_merge.py --base PPT기본양식.pptx --source PPT템플릿_예시.pptx --output 결과.pptx --layout 8

    # 병합 (템플릿 변환 없이)
    python pptx_merge.py --base PPT기본양식.pptx --source PPT템플릿_예시.pptx --output 결과.pptx --layout 8 --no-template
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy
from lxml import etree
import zipfile
import re
import argparse
import sys

# ============================================================
# 설정 - 필요시 수정
# ============================================================

# 제거할 placeholder 텍스트 패턴
REMOVE_PATTERNS = [
    r'우리는 인간생활의 향상과 개선에 필요한[^\n]*',
    r'우리는 제품과 서비스를 생산하기 이전에[^\n]*',
    r'이를 위하여 모든 사람은[^\n]*',
    r'이 목적을 달성하기 위하여[^\n]*',
    r'또한 인재를 양성하고[^\n]*',
    r'우리의 제품과 서비스는[^\n]*',
    r'나아가 문화의 발전에 기여한다[^\n]*',
    r'따라서 기업의 이익은[^\n]*',
]

# 스타일 가이드 -> 템플릿 변환
TEMPLATE_REPLACEMENTS = [
    (r'소제목/?[ ]*Medium[ ]*14pt', '{{소제목}}'),
    (r'중제목[/|]?[ ]*Medium[ ,]*16pt', '{{중제목}}'),
    (r'텍스트를 입력하세요\s*', '{{텍스트}}'),
    (r'텍스트를 입력하시오', '{{텍스트}}'),
]


# ============================================================
# 핵심 함수
# ============================================================

def convert_to_template(text):
    """텍스트를 LLM 템플릿으로 변환"""
    result = text

    # placeholder 문장 제거
    for pattern in REMOVE_PATTERNS:
        result = re.sub(pattern, '', result, flags=re.IGNORECASE)

    # 스타일 가이드 -> 템플릿
    for pattern, replacement in TEMPLATE_REPLACEMENTS:
        result = re.sub(pattern, replacement, result)

    # 소제목 뒤에 본문 템플릿 추가
    if '{{소제목}}' in result and '{{본문_내용}}' not in result:
        result = result.replace('{{소제목}}', '{{소제목}}\n{{본문_내용}}')

    # 정리
    result = re.sub(r'[\s\|\.]+$', '', result)
    result = re.sub(r'^[\s\|\.]+', '', result)
    result = re.sub(r'\n\s*\n+', '\n', result)

    return result.strip()


def merge_pptx(base_pptx, source_pptx, output_pptx, layout_index=None, apply_template=True):
    """
    PPTX 파일 병합

    Args:
        base_pptx: 기본 양식 (슬라이드 마스터 유지)
        source_pptx: 추가할 슬라이드
        output_pptx: 출력 파일
        layout_index: 레이아웃 인덱스 (None이면 빈 레이아웃)
        apply_template: LLM 템플릿 변환 적용 여부
    """
    prs_base = Presentation(base_pptx)
    prs_source = Presentation(source_pptx)

    # 레이아웃 선택
    if layout_index is not None:
        target_layout = prs_base.slide_layouts[layout_index]
    else:
        # 빈 레이아웃 찾기
        try:
            target_layout = prs_base.slide_layouts[6]
        except:
            target_layout = prs_base.slide_layouts[-1]

    print(f"기본 양식: {len(prs_base.slides)}개 슬라이드")
    print(f"추가 대상: {len(prs_source.slides)}개 슬라이드")
    print(f"레이아웃: [{layout_index}] '{target_layout.name}'")
    print(f"템플릿 변환: {'활성화' if apply_template else '비활성화'}")
    print()

    # 슬라이드 복사
    for idx, slide in enumerate(prs_source.slides):
        new_slide = prs_base.slides.add_slide(target_layout)

        # 레이아웃의 기본 placeholder 제거
        for shape in list(new_slide.shapes):
            if shape.is_placeholder:
                shape._element.getparent().remove(shape._element)

        # 원본 shape 복사
        for shape in slide.shapes:
            new_el = deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

        # 템플릿 변환
        if apply_template:
            for shape in new_slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        full_text = ''.join(run.text for run in para.runs)
                        converted = convert_to_template(full_text)
                        if full_text != converted and para.runs:
                            para.runs[0].text = converted
                            for run in para.runs[1:]:
                                run.text = ''

        # 사용된 템플릿 확인
        templates = set(re.findall(r'\{\{[^}]+\}\}',
                        ' '.join(s.text for s in new_slide.shapes if hasattr(s, 'text'))))
        template_str = ', '.join(sorted(templates)) if templates else '-'
        print(f"슬라이드 {idx + 1}: {template_str}")

    prs_base.save(output_pptx)
    print(f"\n✓ 완료: {output_pptx} (총 {len(prs_base.slides)}개 슬라이드)")

    return prs_base


def analyze_pptx(pptx_file):
    """PPTX 파일 분석"""
    prs = Presentation(pptx_file)

    print("=" * 60)
    print(f"파일: {pptx_file}")
    print("=" * 60)
    print(f"\n슬라이드 크기: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print(f"슬라이드 수: {len(prs.slides)}")
    print(f"슬라이드 마스터 수: {len(prs.slide_masters)}")

    # 레이아웃 목록
    for master in prs.slide_masters:
        print(f"\n슬라이드 레이아웃 ({len(master.slide_layouts)}개):")
        print("-" * 50)
        for idx, layout in enumerate(master.slide_layouts):
            ph_count = len(list(layout.placeholders))
            print(f"  [{idx:2d}] '{layout.name}' (placeholder: {ph_count}개)")

    return prs


def analyze_layout(pptx_file, layout_index):
    """특정 레이아웃의 placeholder 상세 분석"""
    prs = Presentation(pptx_file)
    layout = prs.slide_layouts[layout_index]

    print("=" * 60)
    print(f"레이아웃 [{layout_index}]: '{layout.name}'")
    print("=" * 60)

    placeholders = []
    for ph in layout.placeholders:
        placeholders.append({
            'idx': ph.placeholder_format.idx,
            'type': str(ph.placeholder_format.type).replace("PLACEHOLDER_", ""),
            'top': ph.top.inches,
            'left': ph.left.inches,
            'width': ph.width.inches,
            'height': ph.height.inches,
            'text': ph.text_frame.text[:40] if hasattr(ph, 'text_frame') and ph.text_frame.text else ""
        })

    # 위치(top) 기준 정렬
    placeholders.sort(key=lambda x: x['top'])

    print(f"\nPlaceholder ({len(placeholders)}개):")
    print("-" * 50)
    for ph in placeholders:
        print(f"\n[{ph['idx']:2d}] {ph['type']}")
        print(f"    위치: Y={ph['top']:.2f}\", X={ph['left']:.2f}\"")
        print(f"    크기: {ph['width']:.2f}\" x {ph['height']:.2f}\"")
        if ph['text']:
            print(f"    기본텍스트: '{ph['text']}'")


def analyze_theme(pptx_file):
    """테마 색상 및 폰트 분석"""
    print("=" * 60)
    print("테마 분석")
    print("=" * 60)

    with zipfile.ZipFile(pptx_file, 'r') as z:
        theme_files = [f for f in z.namelist() if 'theme1.xml' in f]

        for theme_file in theme_files:
            theme_xml = z.read(theme_file)
            root = etree.fromstring(theme_xml)

            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

            # 색상 스키마
            clrScheme = root.find('.//a:clrScheme', nsmap)
            if clrScheme is not None:
                print(f"\n[색상 스키마: {clrScheme.get('name', 'Unknown')}]")

                color_names = {
                    'dk1': '어두운색1 (텍스트)', 'lt1': '밝은색1 (배경)',
                    'dk2': '어두운색2', 'lt2': '밝은색2',
                    'accent1': '강조색1', 'accent2': '강조색2',
                    'accent3': '강조색3', 'accent4': '강조색4',
                    'accent5': '강조색5', 'accent6': '강조색6',
                    'hlink': '하이퍼링크', 'folHlink': '방문한 링크',
                }

                for elem in clrScheme:
                    tag = elem.tag.split('}')[-1]
                    if tag in color_names:
                        srgb = elem.find('.//a:srgbClr', nsmap)
                        sysclr = elem.find('.//a:sysClr', nsmap)
                        if srgb is not None:
                            print(f"  {color_names[tag]}: #{srgb.get('val')}")
                        elif sysclr is not None:
                            print(f"  {color_names[tag]}: #{sysclr.get('lastClr', 'N/A')}")

            # 폰트 스키마
            fontScheme = root.find('.//a:fontScheme', nsmap)
            if fontScheme is not None:
                print(f"\n[폰트 스키마: {fontScheme.get('name', 'Unknown')}]")

                for font_type, label in [('majorFont', '제목'), ('minorFont', '본문')]:
                    font = fontScheme.find(f'.//a:{font_type}', nsmap)
                    if font is not None:
                        latin = font.find('a:latin', nsmap)
                        ea = font.find('a:ea', nsmap)
                        print(f"  {label}:")
                        if latin is not None:
                            print(f"    - 라틴: {latin.get('typeface')}")
                        if ea is not None and ea.get('typeface'):
                            print(f"    - 동아시아: {ea.get('typeface')}")


# ============================================================
# CLI
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description='PPTX 병합 및 LLM 템플릿 변환',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
예시:
  # 파일 분석
  python pptx_merge.py --analyze PPT기본양식.pptx

  # 레이아웃 상세 분석
  python pptx_merge.py --analyze PPT기본양식.pptx --layout 8

  # 테마 분석
  python pptx_merge.py --analyze PPT기본양식.pptx --theme

  # 병합 (템플릿 변환 포함)
  python pptx_merge.py --base PPT기본양식.pptx --source PPT템플릿_예시.pptx --output 결과.pptx --layout 8

  # 병합 (템플릿 변환 없이)
  python pptx_merge.py --base PPT기본양식.pptx --source PPT템플릿_예시.pptx --output 결과.pptx --no-template
        """
    )

    # 분석 모드
    parser.add_argument('--analyze', metavar='FILE', help='PPTX 파일 분석')
    parser.add_argument('--theme', action='store_true', help='테마 색상/폰트 분석 (--analyze와 함께 사용)')

    # 병합 모드
    parser.add_argument('--base', help='기본 양식 파일 (슬라이드 마스터 유지)')
    parser.add_argument('--source', help='추가할 슬라이드 파일')
    parser.add_argument('--output', help='출력 파일')
    parser.add_argument('--layout', type=int, default=None, help='레이아웃 인덱스')
    parser.add_argument('--no-template', action='store_true', help='템플릿 변환 비활성화')

    args = parser.parse_args()

    # 분석 모드
    if args.analyze:
        analyze_pptx(args.analyze)
        if args.layout is not None:
            print()
            analyze_layout(args.analyze, args.layout)
        if args.theme:
            print()
            analyze_theme(args.analyze)
        return

    # 병합 모드
    if args.base and args.source and args.output:
        merge_pptx(
            args.base,
            args.source,
            args.output,
            args.layout,
            not args.no_template
        )
        return

    # 인자 부족
    parser.print_help()
    sys.exit(1)


if __name__ == "__main__":
    main()
