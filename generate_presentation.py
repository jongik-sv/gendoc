import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def duplicate_slide(pres, index):
    """
    Duplicate a slide at the given index and append it to the end of the presentation.
    """
    source_slide = pres.slides[index]
    source_slide = pres.slides[index]
    dest_slide = pres.slides.add_slide(source_slide.slide_layout)

    # Copy shapes
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return dest_slide

def replace_text(slide, replacements):
    """
    Replace text in a slide based on a dictionary of {placeholder_text: new_text}.
    Also handles simple "Title" and "Body" replacements by looking for text frames.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, value in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, value)
    
    # Heuristic for Title and Body if specific placeholders aren't found
    # This is a simplification; in a real scenario, we might target specific shape indices or names.
    # For this script, we'll assume the template has generic text we want to replace.
    pass

def set_text_by_index(slide, shape_index, text, font_size=None, bold=False):
    """
    Set text for a specific shape index.
    """
    try:
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return
        
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        
        if font_size:
            run.font.size = Pt(font_size)
        if bold:
            run.font.bold = True
            
    except IndexError:
        print(f"Warning: Shape index {shape_index} not found on slide.")

def create_presentation():
    base_pptx = "PPT기본양식_병합.pptx"
    output_pptx = "jjiban_presentation.pptx"
    
    prs = Presentation(base_pptx)
    
    # Template Indices (0-based) based on analysis
    # Title: 0
    # Vision: 10 (Slide 11)
    # Target: 13 (Slide 14)
    # Features: 18 (Slide 19)
    # LLM: 15 (Slide 16)
    # Summary: 37 (Slide 38)
    # Tech Stack: 29 (Slide 30)
    # Architecture: 10 (Slide 11 - reusing generic text slide for custom drawing)
    # Data: 49 (Slide 50)
    # Workflow: 36 (Slide 37)
    # Roadmap: 45 (Slide 46)
    # Q&A: 4 (Slide 5)

    slides_to_create = []

    # 1. Title Slide
    slide = duplicate_slide(prs, 0)
    # Assuming shape 0 is Title, 1 is Subtitle (need to verify, but usually title is early)
    # We will iterate shapes to find text matches or just overwrite specific ones if we knew IDs.
    # Since we don't know exact IDs, we'll use a search-and-replace approach for the template text
    # or just add new text boxes over them if needed. 
    # BUT, let's try to find the title shape.
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Heuristic: Title usually has larger font or specific text
            if "Title" in shape.text_frame.text or "제목" in shape.text_frame.text:
                shape.text_frame.text = "jjiban - AI와 함께하는\n차세대 프로젝트 관리 도구"
            if "Subtitle" in shape.text_frame.text or "부제" in shape.text_frame.text:
                shape.text_frame.text = "개발자 친화적 로컬 기반 PM 도구\n발표자: [발표자 성명]"

    # 2. Vision (Slide 11 template)
    slide = duplicate_slide(prs, 10)
    # Replace title
    for shape in slide.shapes:
        if shape.has_text_frame and ("Title" in shape.text_frame.text or "제목" in shape.text_frame.text):
             shape.text_frame.text = "제품 비전"
             break
    # Add content
    # This template likely has a body text area.
    # We'll add a text box for content if we can't easily find the placeholder.
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = textbox.text_frame
    tf.text = "비전: \"LLM과 함께 개발하는 차세대 프로젝트 관리 도구\""
    p = tf.add_paragraph()
    p.text = "핵심 가치:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Local First: 내 컴퓨터에서 npx jjiban으로 즉시 실행"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Git Friendly: 모든 데이터는 파일로 저장, Git으로 동기화"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• AI Native: LLM이 직접 프로젝트 관리 데이터 수정 가능"
    p.level = 1

    # 3. Target & Problem (Slide 14 template - 3 columns)
    slide = duplicate_slide(prs, 13)
    # Title
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "타겟 사용자 및 해결 과제"
    # Columns content - manually positioning text boxes over the columns
    # Col 1: Target
    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(3), Inches(3))
    tb1.text_frame.text = "타겟 사용자\n\n1~10인 규모의\n소규모 개발팀"
    # Col 2: Problem
    tb2 = slide.shapes.add_textbox(Inches(4), Inches(2.5), Inches(3), Inches(3))
    tb2.text_frame.text = "현재의 문제점\n\n• 무거운 PM 도구\n• AI 도구와 단절\n• 데이터 주권 우려"
    # Col 3: Solution
    tb3 = slide.shapes.add_textbox(Inches(7.5), Inches(2.5), Inches(3), Inches(3))
    tb3.text_frame.text = "jjiban의 해결책\n\n• 설치 없는 실행\n• IDE/Terminal 통합\n• 텍스트 기반 데이터"

    # 4. Key Features (Slide 19 template - 4 icons)
    slide = duplicate_slide(prs, 18)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "핵심 특징"
    # We can overlay text on the 4 quadrants
    # Q1
    slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(2), Inches(1)).text_frame.text = "No Database\n파일 기반 저장"
    # Q2
    slide.shapes.add_textbox(Inches(3.5), Inches(3.5), Inches(2), Inches(1)).text_frame.text = "Git Sync\n팀 동기화"
    # Q3
    slide.shapes.add_textbox(Inches(6), Inches(3.5), Inches(2), Inches(1)).text_frame.text = "Conflict Free\n분산 JSON"
    # Q4
    slide.shapes.add_textbox(Inches(8.5), Inches(3.5), Inches(2), Inches(1)).text_frame.text = "Offline\n로컬 작업"

    # 5. LLM Integration (Slide 16 template)
    slide = duplicate_slide(prs, 15)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "LLM 협업 (AI Integration)"
    # Left: Keyword
    slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(2)).text_frame.text = "CLI 통합\n& 직접 제어"
    # Right: Detail
    slide.shapes.add_textbox(Inches(4.5), Inches(3), Inches(5), Inches(3)).text_frame.text = "• Claude Code, Gemini CLI 연동\n• 파일 시스템을 통한 Task 상태 변경\n• \"이 버그 수정했어\" → JSON 자동 수정\n• 프로젝트 문맥(Context) 완벽 유지"

    # 6. Summary (Slide 38 template - 4 cards)
    slide = duplicate_slide(prs, 37)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "주요 기능 요약"
    # Card 1
    slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(2), Inches(2)).text_frame.text = "WBS 트리 뷰\n계층형 작업 관리"
    # Card 2
    slide.shapes.add_textbox(Inches(3), Inches(4), Inches(2), Inches(2)).text_frame.text = "칸반 보드\n직관적 상태 관리"
    # Card 3
    slide.shapes.add_textbox(Inches(5.5), Inches(4), Inches(2), Inches(2)).text_frame.text = "워크플로우 엔진\n유연한 규칙 정의"
    # Card 4
    slide.shapes.add_textbox(Inches(8), Inches(4), Inches(2), Inches(2)).text_frame.text = "문서 관리\nMarkdown 통합"

    # 7. Tech Stack (Slide 30 template - Table)
    slide = duplicate_slide(prs, 29)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "기술 스택 (Tech Stack)"
    # Overlay table content
    # Assuming table is centrally located, we might just add text boxes over cells if we can't access table easily
    # Or draw a new table. Let's draw a simple text representation for robustness.
    tb = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(3))
    tf = tb.text_frame
    tf.text = "Runtime: Node.js 20.x"
    tf.add_paragraph().text = "Framework: Nuxt 3 (Standalone)"
    tf.add_paragraph().text = "Frontend: Vue 3 + PrimeVue + TailwindCSS"
    tf.add_paragraph().text = "Data: File System (JSON + Markdown)"

    # 8. Architecture (Custom Drawing)
    slide = duplicate_slide(prs, 10) # Blank-ish slide
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "시스템 구조 (System Architecture)"
    
    # Draw Architecture Diagram
    # User (Browser) <-> Nuxt Server <-> File System <-> Git
    
    # User
    user_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(1.5), Inches(1))
    user_shape.text_frame.text = "User\n(Browser)"
    
    # Nuxt Server
    server_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3), Inches(1.5), Inches(1))
    server_shape.text_frame.text = "Nuxt Server\n(Localhost)"
    
    # File System
    fs_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(3), Inches(1.5), Inches(1))
    fs_shape.text_frame.text = "File System\n(.jjiban/)"
    
    # Git
    git_shape = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(8.5), Inches(3), Inches(1.5), Inches(1))
    git_shape.text_frame.text = "Git\n(Remote)"
    
    # Arrows
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.5), Inches(3.4), Inches(1), Inches(0.2))
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5), Inches(3.4), Inches(1), Inches(0.2))
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.5), Inches(3.4), Inches(1), Inches(0.2))

    # 9. Data Structure (Slide 50 template - Hierarchy)
    slide = duplicate_slide(prs, 49)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "데이터 구조 (Data Structure)"
    # Overlay hierarchy text
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = tb.text_frame
    tf.text = ".jjiban/"
    p = tf.add_paragraph()
    p.text = "  ├── projects/ (프로젝트 데이터)"
    p = tf.add_paragraph()
    p.text = "  ├── settings/ (전역 설정)"
    p = tf.add_paragraph()
    p.text = "  └── templates/ (문서 템플릿)"
    p = tf.add_paragraph()
    p.text = "\n핵심 파일:"
    p = tf.add_paragraph()
    p.text = "  • wbs.md: 통합 구조 문서"
    p = tf.add_paragraph()
    p.text = "  • JSON: 개별 Task 상세"

    # 10. Workflow (Slide 37 template - Process Chain)
    slide = duplicate_slide(prs, 36)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "워크플로우 엔진"
    # Overlay text on circles
    # Circle 1
    slide.shapes.add_textbox(Inches(1), Inches(4), Inches(1.5), Inches(1)).text_frame.text = "유연성\n(JSON 정의)"
    # Circle 2
    slide.shapes.add_textbox(Inches(3.5), Inches(4), Inches(1.5), Inches(1)).text_frame.text = "자동화\n(템플릿 생성)"
    # Circle 3
    slide.shapes.add_textbox(Inches(6), Inches(4), Inches(1.5), Inches(1)).text_frame.text = "확장성\n(카테고리별)"
    # Circle 4
    slide.shapes.add_textbox(Inches(8.5), Inches(4), Inches(1.5), Inches(1)).text_frame.text = "통합\n(LLM 연동)"

    # 11. Roadmap (Slide 46 template - Timeline)
    slide = duplicate_slide(prs, 45)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = "향후 계획 (Roadmap)"
    # Overlay timeline events
    slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(2)).text_frame.text = "현재 (PoC)\n• 핵심 기능\n• LLM 검증"
    slide.shapes.add_textbox(Inches(4), Inches(3), Inches(2), Inches(2)).text_frame.text = "Beta\n• Gantt 차트\n• 웹 터미널"
    slide.shapes.add_textbox(Inches(7), Inches(3), Inches(2), Inches(2)).text_frame.text = "v1.0\n• 플러그인\n• LLM 공식 지원"

    # 12. Q&A (Slide 5 template)
    slide = duplicate_slide(prs, 4)
    slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(2)).text_frame.text = "Q & A\n질의응답"


    # Remove original template slides (0-107)
    # We added 12 slides. They are now at indices 108 to 119.
    # We need to delete 0 to 107.
    # Deleting slides in python-pptx is tricky. The standard way is to remove from xml.
    # But since we appended, we can just keep the last 12.
    
    # A safer way to "delete" is to create a NEW presentation and copy the NEW slides over,
    # but copying between presentations is hard.
    # Instead, we will delete the first 108 slides.
    # We must delete from the end of the range backwards to avoid index shifting issues.
    
    xml_slides = prs.slides._sldIdLst
    slides_to_keep = 12
    total_slides = len(prs.slides)
    
    # We want to keep the last `slides_to_keep` slides.
    # So we delete indices 0 to (total_slides - slides_to_keep - 1).
    # We delete in reverse order.
    
    for i in range(total_slides - slides_to_keep - 1, -1, -1):
        prs.slides._sldIdLst.remove(xml_slides[i])

    prs.save(output_pptx)
    print(f"Presentation saved to {output_pptx}")

if __name__ == "__main__":
    create_presentation()
